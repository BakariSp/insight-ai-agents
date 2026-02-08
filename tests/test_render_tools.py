"""Tests for render_tools — PPT outline proposal, generation, and styling."""

import pytest
from pathlib import Path


# ── propose_pptx_outline ─────────────────────────────────────────


@pytest.mark.asyncio
async def test_propose_pptx_outline_returns_structure():
    from tools.render_tools import propose_pptx_outline

    outline_data = [
        {
            "title": "Introduction",
            "key_points": ["Welcome", "Overview", "Objectives"],
            "layout": "title",
            "notes_summary": "Welcome the class",
            "section": "Introduction",
        },
        {
            "title": "Core Concepts",
            "key_points": ["Point A", "Point B", "Point C", "Point D"],
            "layout": "content",
            "notes_summary": "Explain core concepts",
            "section": "Core Content",
        },
    ]

    result = await propose_pptx_outline(
        title="Test Presentation",
        outline=outline_data,
        total_slides=10,
        estimated_duration=45,
    )

    assert result["title"] == "Test Presentation"
    assert result["totalSlides"] == 10
    assert result["estimatedDuration"] == 45
    assert result["status"] == "proposed"
    assert len(result["outline"]) == 2
    assert result["outline"][0]["title"] == "Introduction"


@pytest.mark.asyncio
async def test_propose_pptx_outline_auto_counts_slides():
    from tools.render_tools import propose_pptx_outline

    outline_data = [{"title": f"Slide {i}"} for i in range(15)]

    result = await propose_pptx_outline(
        title="Auto Count",
        outline=outline_data,
    )

    assert result["totalSlides"] == 15
    assert result["estimatedDuration"] == 0


@pytest.mark.asyncio
async def test_propose_pptx_outline_enforces_max_slides():
    from tools.render_tools import propose_pptx_outline
    from config.settings import get_settings

    max_slides = get_settings().pptx_max_slides
    outline_data = [{"title": f"Slide {i}"} for i in range(max_slides + 10)]

    result = await propose_pptx_outline(
        title="Too Many Slides",
        outline=outline_data,
    )

    assert result["totalSlides"] == max_slides
    assert len(result["outline"]) == max_slides


@pytest.mark.asyncio
async def test_propose_pptx_outline_minimal_entries():
    """Outline entries only need a title — all other fields are optional."""
    from tools.render_tools import propose_pptx_outline

    result = await propose_pptx_outline(
        title="Minimal",
        outline=[{"title": "Only Title"}],
    )

    assert result["totalSlides"] == 1
    assert result["outline"][0]["title"] == "Only Title"


# ── generate_pptx ────────────────────────────────────────────────


@pytest.mark.asyncio
async def test_generate_pptx_basic():
    from tools.render_tools import generate_pptx

    slides = [
        {"layout": "title", "title": "Test Title", "body": "Subtitle"},
        {"layout": "content", "title": "Slide 2", "body": "Point 1\nPoint 2\nPoint 3"},
    ]

    result = await generate_pptx(slides=slides, title="Test")

    assert result["filename"] == "Test.pptx"
    assert result["slide_count"] == 2
    assert result["size"] > 0
    assert result["url"]  # non-empty URL


@pytest.mark.asyncio
async def test_generate_pptx_section_header_layout():
    from tools.render_tools import generate_pptx

    slides = [
        {"layout": "title", "title": "Main Title", "body": "Subtitle"},
        {"layout": "section_header", "title": "Part I", "body": "Introduction section"},
        {"layout": "content", "title": "Details", "body": "A\nB\nC"},
    ]

    result = await generate_pptx(slides=slides, title="Section Test")

    assert result["slide_count"] == 3
    assert result["size"] > 0


@pytest.mark.asyncio
async def test_generate_pptx_two_column_layout():
    from tools.render_tools import generate_pptx

    slides = [
        {
            "layout": "two_column",
            "title": "Comparison",
            "left": "Left A\nLeft B",
            "right": "Right A\nRight B",
        },
    ]

    result = await generate_pptx(slides=slides, title="Column Test")

    assert result["slide_count"] == 1
    assert result["size"] > 0


@pytest.mark.asyncio
async def test_generate_pptx_speaker_notes():
    from tools.render_tools import generate_pptx
    from pptx import Presentation
    import tempfile

    slides = [
        {
            "layout": "content",
            "title": "With Notes",
            "body": "Content here",
            "notes": "These are speaker notes for the teacher",
        },
    ]

    result = await generate_pptx(slides=slides, title="Notes Test")

    # Extract the temp filename from the fallback URL and read the actual file
    url = result["url"]
    temp_filename = url.split("/")[-1]
    filepath = Path(tempfile.gettempdir()) / temp_filename
    assert filepath.exists(), f"Generated file not found at {filepath}"

    prs = Presentation(str(filepath))
    notes = prs.slides[0].notes_slide.notes_text_frame.text
    assert "speaker notes for the teacher" in notes


@pytest.mark.asyncio
async def test_generate_pptx_many_slides():
    from tools.render_tools import generate_pptx

    slides = [{"layout": "title", "title": "Main Title", "body": "Overview"}]
    for i in range(1, 16):
        slides.append({
            "layout": "content",
            "title": f"Slide {i}: Topic {i}",
            "body": "\n".join(f"Point {j}" for j in range(1, 6)),
            "notes": f"Notes for slide {i}",
        })

    result = await generate_pptx(slides=slides, title="Many Slides")

    assert result["slide_count"] == 16
    assert result["size"] > 0


@pytest.mark.asyncio
async def test_generate_pptx_education_fallback():
    """When no template file exists, programmatic styling is applied."""
    from tools.render_tools import generate_pptx

    slides = [
        {"layout": "title", "title": "Education Theme", "body": "Auto-styled"},
        {"layout": "content", "title": "Content", "body": "A\nB\nC"},
    ]

    # template="education" is the default — no .pptx file, uses programmatic fallback
    result = await generate_pptx(slides=slides, title="Theme Test", template="education")

    assert result["slide_count"] == 2
    assert result["size"] > 0


# ── SSE Event for outline ────────────────────────────────────────


def test_pptx_outline_sse_event():
    from services.datastream import DataStreamEncoder
    from api.conversation import _build_tool_result_events

    enc = DataStreamEncoder()
    result = {
        "title": "Test PPT",
        "outline": [{"title": "Slide 1"}],
        "totalSlides": 5,
        "estimatedDuration": 30,
        "status": "proposed",
    }

    events = _build_tool_result_events(enc, "propose_pptx_outline", result)

    assert len(events) == 1
    assert "pptx-outline" in events[0]
    assert "Test PPT" in events[0]
    assert "requiresConfirmation" in events[0]


def test_file_ready_sse_event_for_pptx():
    from services.datastream import DataStreamEncoder
    from api.conversation import _build_tool_result_events

    enc = DataStreamEncoder()
    result = {
        "url": "/api/files/generated/test.pptx",
        "filename": "test.pptx",
        "slide_count": 10,
        "size": 50000,
    }

    events = _build_tool_result_events(enc, "generate_pptx", result)

    assert len(events) == 1
    assert "file-ready" in events[0]
    assert "pptx" in events[0]


# ── Helpers ──────────────────────────────────────────────────────


def test_safe_filename():
    from tools.render_tools import _safe_filename

    assert _safe_filename("Hello World") == "Hello World"
    assert _safe_filename('File<>:"/\\|?*Name') == "FileName"
    assert _safe_filename("") == "untitled"
    assert len(_safe_filename("x" * 200)) <= 100


def test_get_template_path_returns_none_for_missing():
    from tools.render_tools import _get_template_path

    result = _get_template_path("nonexistent_template")
    assert result is None
