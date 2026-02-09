"""Render tools — convert structured data into downloadable files.

These tools contain no AI logic.  AI logic lives in the Agent's tool-use loop
managed by the LLM.  Tools only: accept structured input → render file → upload → return URL.
"""

from __future__ import annotations

import re
import tempfile
import uuid
from collections import OrderedDict
from pathlib import Path

from config.settings import get_settings


# ── PPT Outline Proposal ────────────────────────────────────────


async def propose_pptx_outline(
    title: str,
    outline: list[dict],
    total_slides: int = 0,
    estimated_duration: int = 0,
) -> dict:
    """Propose a PPT outline for teacher review before full generation.

    Call this BEFORE generate_pptx to let the teacher review the structure.
    The frontend will show the outline with a confirm/revise UI.

    Args:
        title: Presentation title.
        outline: List of outline entries. Each entry needs at minimum:
            - title (str): Slide or section title
            Any of these optional fields may be included if useful:
            - key_points (list[str]): Main points for this slide
            - section (str): Section grouping label
            - layout (str): Suggested layout hint
        total_slides: Total number of slides (auto-counted if 0).
        estimated_duration: Estimated presentation duration in minutes.

    Returns:
        {"title": ..., "outline": [...], "totalSlides": N, "estimatedDuration": N, "status": "proposed"}
    """
    settings = get_settings()
    slide_count = total_slides or len(outline)
    if slide_count > settings.pptx_max_slides:
        slide_count = settings.pptx_max_slides
        outline = outline[:settings.pptx_max_slides]

    return {
        "title": title,
        "outline": outline,
        "totalSlides": slide_count,
        "estimatedDuration": estimated_duration,
        "status": "proposed",
    }


# ── PPT Generation ──────────────────────────────────────────────


async def generate_pptx(
    slides: list[dict],
    title: str = "Presentation",
    template: str = "education",
) -> dict:
    """Generate a PowerPoint file from structured slides JSON and return a download link.

    Args:
        slides: List of slide dicts, each with layout/title/body/notes fields.
            layout options: "title", "section_header", "content", "two_column".
        title: Presentation title.
        template: Template name (education/default/minimal).

    Returns:
        {"url": "...", "filename": "xxx.pptx", "slide_count": N, "size": bytes}
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # Load template or create a styled blank presentation
    template_path = _get_template_path(template)
    if template_path and template_path.exists():
        prs = Presentation(str(template_path))
    else:
        prs = Presentation()
        _apply_education_theme(prs)

    for slide_data in slides:
        layout_name = slide_data.get("layout", "content")

        if layout_name == "title":
            slide_layout = prs.slide_layouts[0]  # Title Slide
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data.get("title", "")
            _style_title(slide.shapes.title, Pt(36), RGBColor(0x1E, 0x5A, 0x96))
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_data.get("body", "")
                _style_text_frame(slide.placeholders[1].text_frame, Pt(18), RGBColor(0x55, 0x55, 0x55))

        elif layout_name == "section_header":
            # Section divider slide
            try:
                slide_layout = prs.slide_layouts[2]  # Section Header
            except IndexError:
                slide_layout = prs.slide_layouts[0]  # Fallback to title
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data.get("title", "")
            _style_title(slide.shapes.title, Pt(32), RGBColor(0x1E, 0x5A, 0x96))
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_data.get("body", slide_data.get("subtitle", ""))
                _style_text_frame(slide.placeholders[1].text_frame, Pt(16), RGBColor(0x66, 0x66, 0x66))

        elif layout_name == "two_column":
            slide_layout = prs.slide_layouts[3]  # Two Content
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data.get("title", "")
            _style_title(slide.shapes.title, Pt(28), RGBColor(0x1E, 0x5A, 0x96))
            if len(slide.placeholders) > 1:
                _add_bullet_points(
                    slide.placeholders[1].text_frame,
                    slide_data.get("left", "").split("\n"),
                    Pt(14), RGBColor(0x33, 0x33, 0x33),
                )
            if len(slide.placeholders) > 2:
                _add_bullet_points(
                    slide.placeholders[2].text_frame,
                    slide_data.get("right", "").split("\n"),
                    Pt(14), RGBColor(0x33, 0x33, 0x33),
                )

        else:  # "content" or default
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data.get("title", "")
            _style_title(slide.shapes.title, Pt(28), RGBColor(0x1E, 0x5A, 0x96))
            if len(slide.placeholders) > 1:
                body = slide_data.get("body", "")
                _add_bullet_points(
                    slide.placeholders[1].text_frame,
                    body.split("\n"),
                    Pt(16), RGBColor(0x33, 0x33, 0x33),
                )

        # Speaker notes
        if slide_data.get("notes"):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["notes"]

    display_name = f"{_display_filename(title)}.pptx"
    safe_name = f"{_safe_filename(title)}.pptx"
    filepath = Path(tempfile.gettempdir()) / f"{uuid.uuid4().hex}_{safe_name}"
    prs.save(str(filepath))

    url = await _upload_file(
        filepath, display_name,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    return {
        "url": url,
        "filename": display_name,
        "slide_count": len(slides),
        "size": filepath.stat().st_size,
    }


# ── Word Document Generation ────────────────────────────────────


async def generate_docx(
    content: str,
    title: str = "Document",
    format: str = "plain",
) -> dict:
    """Generate a Word document from Markdown content.

    Args:
        content: Markdown-formatted document content.
        title: Document title.
        format: Template style (plain/lesson_plan/worksheet/report).

    Returns:
        {"url": "...", "filename": "xxx.docx", "size": bytes}
    """
    from docx import Document

    doc = Document()
    doc.add_heading(title, 0)

    for line in content.split("\n"):
        line = line.strip()
        if not line:
            continue
        elif line.startswith("### "):
            doc.add_heading(line[4:], level=3)
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif line.startswith("- ") or line.startswith("* "):
            doc.add_paragraph(line[2:], style="List Bullet")
        elif len(line) > 2 and line[0].isdigit() and line[1] == ".":
            doc.add_paragraph(line[3:].strip(), style="List Number")
        elif line.startswith("> "):
            doc.add_paragraph(line[2:])
        elif line.startswith("---"):
            doc.add_page_break()
        else:
            doc.add_paragraph(line)

    display_name = f"{_display_filename(title)}.docx"
    safe_name = f"{_safe_filename(title)}.docx"
    filepath = Path(tempfile.gettempdir()) / f"{uuid.uuid4().hex}_{safe_name}"
    doc.save(str(filepath))

    url = await _upload_file(
        filepath, display_name,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )

    return {
        "url": url,
        "filename": display_name,
        "size": filepath.stat().st_size,
    }


# ── PDF Rendering ────────────────────────────────────────────────


async def render_pdf(
    html_content: str,
    title: str = "Document",
    css_template: str = "default",
) -> dict:
    """Render HTML content into a PDF file.

    Args:
        html_content: HTML string (can include inline CSS).
        title: File title.
        css_template: CSS template (default/worksheet/report).

    Returns:
        {"url": "...", "filename": "xxx.pdf", "size": bytes}
    """
    try:
        from weasyprint import HTML, CSS

        css = _get_css_template(css_template)

        full_html = (
            f'<!DOCTYPE html><html><head><meta charset="utf-8">'
            f"<title>{title}</title></head>"
            f"<body>{html_content}</body></html>"
        )

        pdf_bytes = HTML(string=full_html).write_pdf(
            stylesheets=[CSS(string=css)]
        )
    except ImportError:
        # weasyprint not installed — fallback to returning HTML as-is
        pdf_bytes = html_content.encode("utf-8")

    display_name = f"{_display_filename(title)}.pdf"
    safe_name = f"{_safe_filename(title)}.pdf"
    filepath = Path(tempfile.gettempdir()) / f"{uuid.uuid4().hex}_{safe_name}"
    filepath.write_bytes(pdf_bytes)

    url = await _upload_file(filepath, display_name, "application/pdf")

    return {
        "url": url,
        "filename": display_name,
        "size": len(pdf_bytes),
    }


# ── Interactive HTML Generation ─────────────────────────────────


async def generate_interactive_html(
    html: str,
    title: str = "Interactive Content",
    description: str = "",
    preferred_height: int | None = None,
) -> dict:
    """Generate an interactive HTML web page that is rendered live in the browser.

    Use this tool when the teacher asks for interactive content, simulations,
    animations, drag-and-drop exercises, visual demos, or any web-based
    learning material.  The HTML is rendered in a sandboxed iframe on the
    platform — it supports full HTML/CSS/JavaScript including Canvas, SVG,
    and libraries loaded via CDN (e.g. p5.js, Chart.js, Three.js, D3.js).

    Args:
        html: Complete HTML document string. Must be self-contained (inline
            CSS + JS, or load libraries from CDN).  Can use <!DOCTYPE html>
            or just a <body> fragment — the platform wraps fragments
            automatically.
        title: Short title shown above the preview frame.
        description: One-line description of what this interactive content does.
        preferred_height: Preferred iframe height in pixels (default ~500).

    Returns:
        {"html": "...", "title": "...", "description": "...", "preferredHeight": N}
    """
    return {
        "html": html,
        "title": title,
        "description": description,
        "preferredHeight": preferred_height or 500,
    }


# ── Interactive Content Planning (Three-Stream) ─────────────────


async def request_interactive_content(
    title: str,
    description: str,
    topics: list[str],
    sections: list[dict],
    grade_level: str = "",
    subject: str = "",
    style: str = "modern",
    include_features: list[str] | None = None,
) -> dict:
    """Plan interactive HTML content for three-stream parallel generation.

    Call this instead of generate_interactive_html for better quality and
    progressive rendering. The actual HTML/CSS/JS will be generated in
    parallel as a follow-up step — the teacher sees content appear gradually.

    Define sections with element IDs so the three generators (HTML, CSS, JS)
    stay consistent with each other.

    Args:
        title: Content title shown above the preview frame.
        description: What this interactive content does (one sentence).
        topics: Key concepts to cover (e.g. ["friction", "Newton's laws"]).
        sections: Page structure. Each section needs:
            - id (str): HTML element ID (e.g. "friction-demo")
            - type (str): section type (text/simulation/quiz/chart/animation)
            - desc (str): What this section contains
        grade_level: Target grade (e.g. "Grade 8").
        subject: Subject area (e.g. "Physics").
        style: Visual style — "modern" (default), "playful", or "scientific".
        include_features: Features to include (animation, drag-drop, quiz, simulation).

    Returns:
        Plan dict that triggers three-stream generation.
    """
    return {
        "status": "planned",
        "title": title,
        "description": description,
        "topics": topics,
        "sections": sections,
        "gradeLevel": grade_level,
        "subject": subject,
        "style": style,
        "includeFeatures": include_features or [],
        "willStream": True,
    }


# ── Internal Helpers ─────────────────────────────────────────────


# ── PPT Styling Helpers ──────────────────────────────────────────


def _get_template_path(template: str) -> Path | None:
    """Resolve a template name to a .pptx file path."""
    assets_dir = Path(__file__).resolve().parent.parent / "assets" / "templates"
    path = assets_dir / f"{template}.pptx"
    return path if path.exists() else None


def _apply_education_theme(prs) -> None:
    """Apply an education-themed style programmatically (fallback when no template)."""
    from pptx.util import Inches

    # 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def _style_title(title_shape, font_size, color) -> None:
    """Apply consistent styling to a slide title."""
    if title_shape is None:
        return
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = font_size
            run.font.color.rgb = color
            run.font.bold = True


def _style_text_frame(tf, font_size, color) -> None:
    """Apply consistent styling to all text in a text frame."""
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = font_size
            run.font.color.rgb = color


def _add_bullet_points(tf, lines: list[str], font_size, color) -> None:
    """Add formatted bullet points to a text frame with proper spacing."""
    from pptx.util import Pt

    tf.clear()
    first = True
    for line in lines:
        line = line.strip()
        if not line:
            continue
        # Strip leading bullet markers (-, *, •)
        clean = line.lstrip("-*• ").strip()
        if not clean:
            continue

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.text = clean
        p.space_after = Pt(6)
        p.space_before = Pt(2)
        for run in p.runs:
            run.font.size = font_size
            run.font.color.rgb = color


def _safe_filename(name: str) -> str:
    """Sanitize a filename to ASCII-safe characters for URL paths.

    Non-ASCII characters (Chinese, etc.) are stripped so that the URL
    ``/api/files/generated/{uuid}_{safe_name}.ext`` is always valid ASCII.
    The original display name with Chinese characters is preserved separately
    via ``_display_filename``.
    """
    clean = re.sub(r'[<>:"/\\|?*]', "", name)
    clean = re.sub(r'[^\x20-\x7E]', "", clean)
    clean = re.sub(r'\s+', '_', clean.strip())
    return clean[:100] or "document"


def _display_filename(name: str) -> str:
    """Sanitize a filename for display — preserves Unicode, removes FS-unsafe chars."""
    clean = re.sub(r'[<>:"/\\|?*]', "", name)
    return clean.strip()[:100] or "untitled"


# Module-level mapping: temp filename → original display filename.
# Used by api/files.py to set Content-Disposition with the Chinese name.
# Lost on restart — acceptable for the dev/fallback endpoint.
_MAX_DISPLAY_NAME_CACHE = 2048
_FILE_DISPLAY_NAMES: OrderedDict[str, str] = OrderedDict()


def remember_display_name(temp_filename: str, display_name: str) -> None:
    """Store display filename in a bounded in-memory cache."""
    if not temp_filename:
        return
    if temp_filename in _FILE_DISPLAY_NAMES:
        _FILE_DISPLAY_NAMES.move_to_end(temp_filename)
    _FILE_DISPLAY_NAMES[temp_filename] = display_name
    while len(_FILE_DISPLAY_NAMES) > _MAX_DISPLAY_NAME_CACHE:
        _FILE_DISPLAY_NAMES.popitem(last=False)


def resolve_display_name(temp_filename: str) -> str | None:
    """Resolve display filename from in-memory cache."""
    return _FILE_DISPLAY_NAMES.get(temp_filename)


def _get_css_template(template: str) -> str:
    """Load a CSS template for PDF rendering."""
    templates = {
        "default": """
            body { font-family: 'Noto Sans SC', sans-serif; padding: 2cm; line-height: 1.6; }
            h1 { color: #1a1a1a; border-bottom: 2px solid #333; padding-bottom: 8px; }
            h2 { color: #333; margin-top: 1.5em; }
            table { border-collapse: collapse; width: 100%; margin: 1em 0; }
            th, td { border: 1px solid #ddd; padding: 8px; text-align: left; }
            th { background: #f5f5f5; }
        """,
        "worksheet": """
            body { font-family: 'Noto Sans SC', sans-serif; padding: 1.5cm; font-size: 14pt; }
            h1 { text-align: center; font-size: 20pt; }
            .question { margin: 1.5em 0; page-break-inside: avoid; }
            .answer-area { border: 1px dashed #ccc; min-height: 3cm; margin: 0.5em 0; }
        """,
        "report": """
            body { font-family: 'Noto Sans SC', sans-serif; padding: 2cm; }
            h1 { color: #1a56db; }
            .chart { text-align: center; margin: 1em 0; }
            .summary { background: #f0f4ff; padding: 1em; border-radius: 8px; }
        """,
    }
    return templates.get(template, templates["default"])


async def _upload_file(filepath: Path, filename: str, content_type: str) -> str:
    """Upload a file to storage.

    Tries Java backend OSS upload first; falls back to local serving path.
    The local path is served by the Python Agent's ``GET /api/files/generated/``
    endpoint and proxied through Next.js.
    """
    try:
        from services.java_client import get_java_client
        java_client = get_java_client()
        # Use the underlying httpx client directly for multipart upload,
        # since the JavaClient wrapper only supports JSON payloads.
        http = java_client._http  # noqa: SLF001
        if http is not None:
            with open(filepath, "rb") as f:
                files = {"file": (filename, f, content_type)}
                response = await http.post(
                    "/studio/teacher/me/files/upload",
                    files=files,
                    params={"purpose": "STUDIO"},
                )
                if response.status_code == 200:
                    data = response.json()
                    url = data.get("data", {}).get("fileUrl", "")
                    if url:
                        return url
    except Exception as exc:
        import logging
        logging.getLogger(__name__).debug("OSS upload failed, using local fallback: %s", exc)

    # Fallback: serve from Python Agent's local file endpoint.
    # Return an absolute URL so the frontend (running on a different port)
    # can resolve the download link without guessing the AI Agent origin.
    remember_display_name(filepath.name, filename)
    settings = get_settings()
    base_url = f"http://localhost:{settings.service_port}"
    return f"{base_url}/api/files/generated/{filepath.name}"
    # Note: temporary file is NOT deleted here — caller or cleanup job handles it
