"""Unified RAG engine wrapping RAG-Anything + LightRAG workspace management.

Each teacher gets an isolated workspace (``teacher-{id}``).
A ``public`` workspace holds system-wide knowledge (curriculum, rubrics).

Phase 1: text documents only (digital PDF, DOCX, PPTX) — no OCR/VLM.
"""

from __future__ import annotations

import logging
import os
import tempfile
from pathlib import Path
from typing import Any
from urllib.parse import urlparse

import httpx
import numpy as np

from config.settings import get_settings

logger = logging.getLogger(__name__)

# Module-level singleton
_engine: InsightRAGEngine | None = None


class InsightRAGEngine:
    """Encapsulates RAG-Anything instances with per-teacher workspace isolation."""

    def __init__(self, pg_uri: str, embedding_model: str, embedding_dim: int) -> None:
        self._pg_uri = pg_uri
        self._embedding_model = embedding_model
        self._embedding_dim = embedding_dim
        self._instances: dict[str, Any] = {}  # workspace_id → RAGAnything
        self._file_registry: dict[str, list[dict[str, str]]] = {}  # workspace_id → [{file_id, file_name}]
        self._initialized = False
        self._pg_pool = None  # asyncpg connection pool

    async def initialize(self) -> None:
        """Create connection pool and verify database connectivity (called once at startup)."""
        if self._initialized:
            return
        try:
            import asyncpg
            self._pg_pool = await asyncpg.create_pool(
                dsn=self._pg_uri,
                min_size=2,
                max_size=10,
                max_inactive_connection_lifetime=300,
            )
            async with self._pg_pool.acquire() as conn:
                await conn.execute("SELECT 1")
            logger.info(
                "RAG engine PostgreSQL pool created (min=2, max=10) — %s",
                self._pg_uri,
            )
        except Exception as exc:
            logger.warning(
                "RAG engine PostgreSQL not available: %s — "
                "document parsing will fail until DB is accessible",
                exc,
            )
        self._initialized = True

    def _ensure_pg_env_vars(self) -> None:
        """Set PostgreSQL env vars from pg_uri — LightRAG reads os.environ directly."""
        parsed = urlparse(self._pg_uri)
        env_map = {
            "POSTGRES_HOST": parsed.hostname or "localhost",
            "POSTGRES_PORT": str(parsed.port or 5432),
            "POSTGRES_USER": parsed.username or "insight",
            "POSTGRES_PASSWORD": parsed.password or "",
            "POSTGRES_DATABASE": parsed.path.lstrip("/") or "insight_agent",
        }
        for key, value in env_map.items():
            if key not in os.environ:
                os.environ[key] = value
                logger.debug("Set %s from pg_uri", key)

    def _build_embedding_func(self) -> Any:
        """Build EmbeddingFunc for LightRAG using DashScope embeddings API."""
        from lightrag.utils import EmbeddingFunc

        model = self._embedding_model  # e.g. "text-embedding-v3"
        dim = self._embedding_dim

        async def _embed(texts: list[str]) -> np.ndarray:
            api_key = os.environ.get("DASHSCOPE_API_KEY", "")
            if not api_key:
                raise RuntimeError(
                    "DASHSCOPE_API_KEY not set — cannot generate embeddings for RAG. "
                    "Set it in .env or environment variables."
                )
            try:
                async with httpx.AsyncClient(timeout=60) as client:
                    resp = await client.post(
                        "https://dashscope.aliyuncs.com/compatible-mode/v1/embeddings",
                        headers={"Authorization": f"Bearer {api_key}"},
                        json={"model": model, "input": texts, "encoding_format": "float"},
                    )
                    resp.raise_for_status()
                data = resp.json()
                embeddings = np.array([item["embedding"] for item in data["data"]])
                logger.debug("Embedded %d texts → shape %s", len(texts), embeddings.shape)
                return embeddings
            except httpx.HTTPStatusError as exc:
                logger.error(
                    "DashScope embedding API error %d: %s",
                    exc.response.status_code, exc.response.text[:200],
                )
                raise
            except Exception as exc:
                logger.error("Embedding function failed: %s", exc)
                raise

        return EmbeddingFunc(
            embedding_dim=dim,
            max_token_size=8192,
            func=_embed,
            model_name=model,
        )

    def _build_llm_func(self):
        """Build LLM function for LightRAG KG extraction using LiteLLM."""
        async def _llm_func(
            prompt: str,
            system_prompt: str | None = None,
            history_messages: list | None = None,
            keyword_extraction: bool = False,
            **kwargs,
        ) -> str:
            import litellm
            from services.concurrency import rate_limited_llm_call

            messages: list[dict] = []
            if system_prompt:
                messages.append({"role": "system", "content": system_prompt})
            if history_messages:
                messages.extend(history_messages)
            messages.append({"role": "user", "content": prompt})

            settings = get_settings()
            task = "keyword_extraction" if keyword_extraction else "kg_extraction"
            try:
                resp = await rate_limited_llm_call(
                    litellm.acompletion,
                    model=settings.default_model,
                    messages=messages,
                    max_tokens=settings.max_tokens,
                    timeout=60,
                )
                content = resp.choices[0].message.content
                if not content:
                    logger.warning("RAG LLM returned empty content for %s", task)
                    return ""
                return content
            except Exception as exc:
                logger.error("RAG LLM call failed (%s): %s", task, exc)
                raise

        return _llm_func

    async def get_instance(self, workspace_id: str) -> Any:
        """Get or create a RAGAnything instance for the given workspace.

        Args:
            workspace_id: e.g. "teacher-123" or "public"

        Returns:
            A RAGAnything instance bound to the workspace.
        """
        if workspace_id in self._instances:
            return self._instances[workspace_id]

        from raganything import RAGAnything
        from lightrag import LightRAG

        # Ensure PostgreSQL env vars are set (LightRAG checks os.environ)
        self._ensure_pg_env_vars()

        # Create working directory — use shared parent; LightRAG's NetworkXStorage
        # creates a {workspace}/ subdirectory inside working_dir automatically.
        working_dir = "./rag_workspaces"
        os.makedirs(working_dir, exist_ok=True)

        lightrag = LightRAG(
            working_dir=working_dir,
            # PG for KV, vector, doc_status; NetworkX for graph (no AGE extension needed)
            graph_storage="NetworkXStorage",
            kv_storage="PGKVStorage",
            vector_storage="PGVectorStorage",
            doc_status_storage="PGDocStatusStorage",
            workspace=workspace_id,
            embedding_func=self._build_embedding_func(),
            llm_model_func=self._build_llm_func(),
            embedding_func_max_async=16,
        )

        # Initialize PostgreSQL storage connections
        await lightrag.initialize_storages()

        rag = RAGAnything(
            lightrag=lightrag,
            # Phase 1: no VLM — text documents only
            vision_model_func=None,
        )

        self._instances[workspace_id] = rag
        logger.info("Created RAGAnything instance for workspace '%s'", workspace_id)
        return rag

    async def ingest_document(
        self,
        teacher_id: str,
        file_path: str,
        file_name: str = "",
        file_id: str = "",
    ) -> dict[str, Any]:
        """Parse a document and ingest into the teacher's workspace.

        Strategy:
          1. Try RAGAnything full pipeline (parser + KG + embedding).
          2. On ANY failure, fall back to plain text extraction + LightRAG.ainsert().

        Args:
            teacher_id: The verified teacher ID.
            file_path: Local path to the downloaded file.
            file_name: Original file name (for metadata).
            file_id: Java backend file ID (for citation tracking).

        Returns:
            {"chunk_count": N, "method": str} on success.
        """
        workspace_id = f"teacher-{teacher_id}"
        rag = await self.get_instance(workspace_id)

        logger.info("Ingesting document '%s' into workspace '%s'", file_name, workspace_id)

        # ── Attempt 1: RAGAnything full pipeline ────────────────────
        try:
            await rag.process_document_complete(
                file_path=file_path,
                file_name=file_name,
            )
            logger.info(
                "RAGAnything full pipeline completed for '%s' in workspace '%s'",
                file_name, workspace_id,
            )
            self._register_file(workspace_id, file_id, file_name)
            return {"chunk_count": -1, "method": "raganything"}
        except Exception as exc:
            logger.warning(
                "RAGAnything full pipeline failed for '%s': %s — "
                "falling back to text extraction + LightRAG.ainsert()",
                file_name, exc,
            )

        # ── Attempt 2: Text extraction + direct LightRAG insert ────
        text = _extract_text(file_path)
        if not text.strip():
            raise RuntimeError(
                f"No text could be extracted from '{file_name}' "
                f"(ext={Path(file_path).suffix})"
            )

        text_len = len(text)
        logger.info(
            "Extracted %d characters from '%s', inserting via LightRAG.ainsert()",
            text_len, file_name,
        )

        await rag.lightrag.ainsert(text, file_paths=[file_name])

        logger.info(
            "LightRAG.ainsert() completed for '%s' (%d chars) in workspace '%s'",
            file_name, text_len, workspace_id,
        )
        self._register_file(workspace_id, file_id, file_name)
        return {"chunk_count": -1, "method": "text_fallback"}

    def _register_file(self, workspace_id: str, file_id: str, file_name: str) -> None:
        """Record a file in the in-memory registry for citation tracking."""
        if not file_id:
            return
        registry = self._file_registry.setdefault(workspace_id, [])
        # Avoid duplicates on re-parse
        registry[:] = [f for f in registry if f["file_id"] != file_id]
        registry.append({"file_id": file_id, "file_name": file_name})
        logger.debug(
            "Registered file %s (%s) in workspace '%s' — %d files total",
            file_id, file_name, workspace_id, len(registry),
        )

    def get_workspace_files(self, teacher_id: str) -> list[dict[str, str]]:
        """Return the list of ingested files for a teacher's workspace."""
        workspace_id = f"teacher-{teacher_id}"
        return list(self._file_registry.get(workspace_id, []))

    async def search(
        self,
        teacher_id: str,
        query: str,
        mode: str = "hybrid",
        include_public: bool = True,
        top_k: int = 5,
    ) -> list[dict[str, Any]]:
        """Search the teacher's workspace (and optionally public workspace).

        Args:
            teacher_id: The verified teacher ID.
            query: Natural-language search query.
            mode: Search mode — "hybrid" (vector + KG), "naive" (vector only), "kg" (graph only).
            include_public: Also search the public workspace.
            top_k: Max number of results.

        Returns:
            List of {"content", "source_file", "score", "metadata"} dicts.
        """
        results: list[dict[str, Any]] = []

        # Map our mode names to RAGAnything mode names
        rag_mode = {"hybrid": "mix", "naive": "naive", "kg": "local"}.get(mode, mode)

        logger.info(
            "RAG search: teacher=%s, query='%s', mode=%s→%s",
            teacher_id, query[:80], mode, rag_mode,
        )

        # Search teacher's workspace
        workspace_id = f"teacher-{teacher_id}"
        try:
            rag = await self.get_instance(workspace_id)
            teacher_results = await rag.aquery(query, mode=rag_mode)
            if teacher_results:
                results.append({
                    "content": str(teacher_results),
                    "source": workspace_id,
                    "score": 1.0,
                })
                logger.info(
                    "RAG search returned %d chars from workspace '%s'",
                    len(str(teacher_results)), workspace_id,
                )
            else:
                logger.warning(
                    "RAG search returned empty from workspace '%s' — "
                    "check that documents were ingested with working LLM+embedding",
                    workspace_id,
                )
        except Exception as exc:
            logger.warning("Search in workspace '%s' failed: %s", workspace_id, exc)

        # Optionally search public workspace
        if include_public:
            try:
                public_rag = await self.get_instance("public")
                public_results = await public_rag.aquery(query, mode=rag_mode)
                if public_results:
                    results.append({
                        "content": str(public_results),
                        "source": "public",
                        "score": 0.8,
                    })
            except Exception as exc:
                logger.debug("Search in public workspace failed (may not exist): %s", exc)

        if not results:
            logger.warning(
                "RAG search returned 0 results for teacher=%s query='%s'. "
                "Possible causes: (1) DASHSCOPE_API_KEY not set, "
                "(2) LLM errors during ingestion left VDB tables empty, "
                "(3) document not yet ingested",
                teacher_id, query[:80],
            )

        return results

    async def delete_document(
        self,
        teacher_id: str,
        file_id: str,
        file_name: str = "",
    ) -> dict[str, Any]:
        """Delete all indexed data for a specific document.

        Strategy:
          1. Look up doc_ids from lightrag_doc_status by file_path (= file_name).
          2. Fall back to in-memory _file_registry if file_name not provided.
          3. Call LightRAG.adelete_by_doc_id() for each matching doc_id.
          4. Clean up _file_registry.

        Args:
            teacher_id: The verified teacher ID.
            file_id: Java backend file ID.
            file_name: Original file name (for matching doc_status records).

        Returns:
            {"deleted_doc_ids": [...], "errors": [...]}.
        """
        workspace_id = f"teacher-{teacher_id}"
        deleted: list[str] = []
        errors: list[str] = []

        # Resolve file_name from registry if not provided
        if not file_name:
            registry = self._file_registry.get(workspace_id, [])
            for entry in registry:
                if entry["file_id"] == file_id:
                    file_name = entry["file_name"]
                    break

        if not file_name:
            logger.warning(
                "Cannot resolve file_name for file_id=%s — "
                "will attempt workspace-wide doc_status scan",
                file_id,
            )

        # Query doc_status for matching doc_ids
        doc_ids = await self._find_doc_ids(workspace_id, file_name)

        if not doc_ids:
            logger.warning(
                "No doc_ids found for file_id=%s, file_name='%s' in workspace '%s'",
                file_id, file_name, workspace_id,
            )
            return {"deleted_doc_ids": [], "errors": ["no matching documents found"]}

        # Get LightRAG instance and delete each doc
        try:
            rag = await self.get_instance(workspace_id)
        except Exception as exc:
            msg = f"Failed to get RAG instance for {workspace_id}: {exc}"
            logger.error(msg)
            return {"deleted_doc_ids": [], "errors": [msg]}

        for doc_id in doc_ids:
            try:
                result = await rag.lightrag.adelete_by_doc_id(
                    doc_id, delete_llm_cache=True,
                )
                deleted.append(doc_id)
                logger.info(
                    "Deleted doc_id=%s from workspace '%s' (result=%s)",
                    doc_id, workspace_id, result,
                )
            except Exception as exc:
                msg = f"Failed to delete doc_id={doc_id}: {exc}"
                logger.error(msg)
                errors.append(msg)

        # Clean up in-memory file registry
        registry = self._file_registry.get(workspace_id, [])
        registry[:] = [f for f in registry if f["file_id"] != file_id]

        logger.info(
            "Document deletion complete: file_id=%s, teacher_id=%s, "
            "deleted=%d, errors=%d",
            file_id, teacher_id, len(deleted), len(errors),
        )
        return {"deleted_doc_ids": deleted, "errors": errors}

    async def _find_doc_ids(
        self, workspace_id: str, file_name: str,
    ) -> list[str]:
        """Look up LightRAG doc_ids from the doc_status table by file_path."""
        if not self._pg_pool:
            logger.warning("No PG pool — cannot look up doc_ids")
            return []

        try:
            async with self._pg_pool.acquire() as conn:
                if file_name:
                    rows = await conn.fetch(
                        "SELECT id FROM lightrag_doc_status "
                        "WHERE workspace = $1 AND file_path = $2",
                        workspace_id, file_name,
                    )
                else:
                    # Without file_name, return all docs in workspace
                    rows = await conn.fetch(
                        "SELECT id FROM lightrag_doc_status "
                        "WHERE workspace = $1",
                        workspace_id,
                    )
                return [r["id"] for r in rows]
        except Exception as exc:
            logger.error("Failed to query doc_status: %s", exc)
            return []

    async def close(self) -> None:
        """Gracefully shut down all RAG instances and the connection pool."""
        count = len(self._instances)
        for ws_id, instance in self._instances.items():
            try:
                if hasattr(instance, 'close'):
                    await instance.close()
            except Exception as exc:
                logger.warning("Error closing RAG instance '%s': %s", ws_id, exc)
        self._instances.clear()

        if self._pg_pool is not None:
            await self._pg_pool.close()
            self._pg_pool = None

        logger.info("RAG engine shut down — %d instances closed", count)


def _extract_text(file_path: str) -> str:
    """Extract plain text from a document without external tools (LibreOffice/MineruParser).

    Supports: PDF, DOCX, PPTX, XLSX, TXT, MD, CSV.
    """
    ext = Path(file_path).suffix.lower()

    if ext == ".docx":
        from docx import Document
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    elif ext == ".pdf":
        import fitz  # PyMuPDF
        with fitz.open(file_path) as pdf:
            return "\n".join(page.get_text() for page in pdf)

    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(file_path)
        lines: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
        return "\n".join(lines)

    elif ext in (".xlsx", ".xls"):
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        lines: list[str] = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                row_text = "\t".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    lines.append(row_text)
        wb.close()
        return "\n".join(lines)

    elif ext in (".txt", ".md", ".csv"):
        return Path(file_path).read_text(encoding="utf-8", errors="ignore")

    else:
        raise RuntimeError(f"Unsupported file type for text extraction: {ext}")


async def download_file(url: str, dest_dir: str | None = None) -> str:
    """Download a file from OSS URL to a local temp path.

    Returns:
        Local file path of the downloaded file.
    """
    dest_dir = dest_dir or tempfile.mkdtemp(prefix="rag_download_")
    async with httpx.AsyncClient(verify=False, timeout=120) as client:
        resp = await client.get(url)
        resp.raise_for_status()
        content = resp.content

    # Extract filename from URL or Content-Disposition
    file_name = url.split("/")[-1].split("?")[0] or "document"
    file_path = os.path.join(dest_dir, file_name)

    with open(file_path, "wb") as f:
        f.write(content)

    logger.info("Downloaded file to %s (%d bytes)", file_path, len(content))
    return file_path


def get_rag_engine() -> InsightRAGEngine:
    """Return the module-level RAG engine singleton."""
    if _engine is None:
        raise RuntimeError("RAG engine not initialized — call init_rag_engine() first")
    return _engine


def init_rag_engine() -> InsightRAGEngine:
    """Create and return the global RAG engine singleton."""
    global _engine
    if _engine is not None:
        return _engine

    settings = get_settings()
    _engine = InsightRAGEngine(
        pg_uri=settings.pg_uri,
        embedding_model=settings.embedding_model,
        embedding_dim=settings.embedding_dim,
    )
    return _engine
