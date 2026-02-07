"""Multimodal helpers — build PydanticAI-compatible user prompts with images & documents.

Converts ``Attachment`` objects from ``ConversationRequest`` into:
- ``ImageUrl`` content parts for images (sent to vision model)
- Extracted text blocks for documents (PDF, DOCX, PPTX, XLSX, TXT, CSV)

When no meaningful attachments are present, returns a plain ``str`` so the
agent call path is zero-cost compatible with the existing text-only flow.
"""

from __future__ import annotations

import logging
import os
import tempfile
from pathlib import Path
from typing import Sequence

from pydantic_ai.messages import ImageUrl, UserContent

from models.conversation import Attachment

logger = logging.getLogger(__name__)

# MIME prefix → category mapping
_IMAGE_PREFIXES = ("image/",)
_DOCUMENT_MIMES: set[str] = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",  # docx
    "application/msword",  # doc
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",  # pptx
    "application/vnd.ms-powerpoint",  # ppt
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",  # xlsx
    "application/vnd.ms-excel",  # xls
    "text/plain",
    "text/csv",
    "text/markdown",
}


def _is_image(att: Attachment) -> bool:
    return any(att.mime_type.startswith(p) for p in _IMAGE_PREFIXES)


def _is_document(att: Attachment) -> bool:
    if att.mime_type in _DOCUMENT_MIMES:
        return True
    # Fallback: check file extension
    ext = Path(att.filename).suffix.lower() if att.filename else ""
    return ext in {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls", ".csv", ".txt", ".md"}


def has_images(attachments: list[Attachment] | None) -> bool:
    """Check whether any attachment is an image."""
    if not attachments:
        return False
    return any(_is_image(a) for a in attachments)


def has_attachments(attachments: list[Attachment] | None) -> bool:
    """Check whether any meaningful attachment (image or document) is present."""
    if not attachments:
        return False
    return any(_is_image(a) or _is_document(a) for a in attachments)


async def build_user_content(
    text_prompt: str,
    attachments: list[Attachment],
) -> str | Sequence[UserContent]:
    """Build a PydanticAI-compatible user prompt.

    - No attachments → returns the original ``str`` (zero overhead).
    - With image attachments → refreshes signed URLs and returns
      ``list[UserContent]`` containing ``ImageUrl`` + text.
    - With document attachments → downloads, extracts text, and prepends
      as context to the text prompt.

    Args:
        text_prompt: The assembled text prompt (language hint + history + message).
        attachments: Attachments from the conversation request.

    Returns:
        Plain string or a sequence of ``UserContent`` items.
    """
    image_atts = [a for a in attachments if _is_image(a)]
    doc_atts = [a for a in attachments if _is_document(a)]

    if not image_atts and not doc_atts:
        return text_prompt

    # ── Extract text from document attachments ──
    doc_context = ""
    if doc_atts:
        doc_texts: list[str] = []
        for att in doc_atts:
            try:
                extracted = await _extract_document_text(att)
                if extracted.strip():
                    doc_texts.append(f"[Attached file: {att.filename}]\n{extracted}")
                    logger.info(
                        "Extracted %d chars from document %s (%s)",
                        len(extracted), att.filename, att.mime_type,
                    )
            except Exception as exc:
                logger.warning(
                    "Failed to extract text from %s (%s): %s",
                    att.filename, att.mime_type, exc,
                )
                doc_texts.append(f"[Attached file: {att.filename} — could not extract content]")

        if doc_texts:
            doc_context = "\n\n".join(doc_texts) + "\n\n"

    enriched_prompt = doc_context + text_prompt if doc_context else text_prompt

    # ── If no images, return enriched text directly ──
    if not image_atts:
        logger.info(
            "Built document-enriched prompt: %d doc(s), text %d chars",
            len(doc_atts), len(enriched_prompt),
        )
        return enriched_prompt

    # ── Build multimodal parts with images ──
    parts: list[UserContent] = []

    for att in image_atts:
        fresh_url = await _refresh_url(att.file_id, att.url)
        parts.append(ImageUrl(url=fresh_url))
        logger.debug("Multimodal: added image %s (%s)", att.file_id, att.mime_type)

    parts.append(enriched_prompt)

    logger.info(
        "Built multimodal prompt: %d image(s) + %d doc(s) + text (%d chars)",
        len(image_atts), len(doc_atts), len(enriched_prompt),
    )
    return parts


# ── Document text extraction ──────────────────────────────────────


async def _extract_document_text(att: Attachment) -> str:
    """Download a document attachment and extract its text content."""
    fresh_url = await _refresh_url(att.file_id, att.url)
    file_path = await _download_file(fresh_url, att.filename)
    try:
        return _extract_text(file_path)
    finally:
        # Clean up temp file
        try:
            os.unlink(file_path)
        except OSError:
            pass


async def _download_file(url: str, filename: str = "") -> str:
    """Download a file from a signed URL to a temp path."""
    import httpx

    async with httpx.AsyncClient(verify=False, timeout=120) as client:
        resp = await client.get(url)
        resp.raise_for_status()

    suffix = Path(filename).suffix if filename else ""
    fd, file_path = tempfile.mkstemp(suffix=suffix, prefix="chat_att_")
    os.close(fd)

    with open(file_path, "wb") as f:
        f.write(resp.content)

    logger.debug("Downloaded attachment to %s (%d bytes)", file_path, len(resp.content))
    return file_path


def _extract_text(file_path: str) -> str:
    """Extract plain text from a document file.

    Supports: PDF, DOCX, PPTX, XLSX/XLS/CSV, TXT/MD.
    """
    ext = Path(file_path).suffix.lower()

    if ext == ".pdf":
        import fitz  # PyMuPDF

        with fitz.open(file_path) as pdf:
            return "\n".join(page.get_text() for page in pdf)

    if ext in (".docx",):
        from docx import Document

        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    if ext in (".pptx",):
        from pptx import Presentation

        prs = Presentation(file_path)
        texts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
            if slide_texts:
                texts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
        return "\n\n".join(texts)

    if ext in (".xlsx", ".xls"):
        import openpyxl

        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        sheets: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cell_values = [str(c) if c is not None else "" for c in row]
                if any(v.strip() for v in cell_values):
                    rows.append("\t".join(cell_values))
            if rows:
                sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(sheets)

    if ext == ".csv":
        return Path(file_path).read_text(encoding="utf-8", errors="ignore")

    if ext in (".txt", ".md"):
        return Path(file_path).read_text(encoding="utf-8", errors="ignore")

    raise RuntimeError(f"Unsupported file type for text extraction: {ext}")


# ── URL helpers ──────────────────────────────────────────────────


async def _refresh_url(file_id: str, fallback_url: str) -> str:
    """Get a fresh signed URL from Java backend, falling back to the provided URL."""
    from services.java_file_client import get_file_url

    try:
        url = await get_file_url(file_id)
        if url:
            return url
    except Exception as exc:
        logger.warning("Failed to refresh URL for %s: %s; using fallback", file_id, exc)

    return fallback_url
